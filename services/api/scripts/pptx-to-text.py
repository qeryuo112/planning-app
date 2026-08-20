#!/usr/bin/env python3
"""
从 PPTX 文件提取每页文本内容。
用法: python pptx-to-text.py <input_pptx_path>
输出: stdout 打印 JSON {"slides": [{"pageNumber": 1, "text": "..."}, ...]}
"""
import sys
import json
import os


def main():
    if len(sys.argv) < 2:
        print("Usage: pptx-to-text.py <input_pptx_path>", file=sys.stderr)
        sys.exit(1)

    input_path = sys.argv[1]

    if not os.path.exists(input_path):
        print(json.dumps({"error": f"File not found: {input_path}"}), file=sys.stderr)
        sys.exit(1)

    try:
        from pptx import Presentation
    except ImportError:
        print(json.dumps({"error": "python-pptx not installed"}), file=sys.stderr)
        sys.exit(1)

    prs = Presentation(input_path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slide_text = "\n".join(texts)
        slides.append({"pageNumber": i, "text": slide_text})

    print(json.dumps({"slides": slides}))


if __name__ == "__main__":
    main()
